"""PPTX file parser using python-pptx."""

from __future__ import annotations

from chunky.parsers.base import BaseParser


class PptxParser(BaseParser):
    SUPPORTED_EXTENSIONS = {".pptx", ".ppt"}

    def supports(self, file_path: str) -> bool:
        return any(file_path.lower().endswith(ext) for ext in self.SUPPORTED_EXTENSIONS)

    def parse(self, file_path: str) -> str:
        from pptx import Presentation

        prs = Presentation(file_path)
        slides_text: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = [f"--- Slide {slide_num} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            parts.append(text)
            slides_text.append("\n".join(parts))
        return "\n\n".join(slides_text)
